#!/usr/bin/env python3
"""PowerPoint テンプレートへ物件情報と画像を差し込む（パイプライン⑤）。

    python scripts/pptx_fill.py \
        --template templates/mysoku_a4.pptx \
        --data out.json \
        --images-dir ./attachments \
        --out 資料_渋谷ハイツ.pptx

テンプレート側の約束ごとは 2 つだけ:

  1. テキストは ``{{property_name}}`` のように二重波かっこで書く。
     キーは assets/property_fields.json の ``pptx`` の値。
  2. 画像を入れたい場所には図形（四角形など）を置き、その図形の**名前**を
     ``IMG:exterior`` / ``IMG:floor_plan`` / ``IMG:map`` / ``IMG:interior``
     にする。図形の位置とサイズが画像の枠になり、図形自体は削除される。
     図形名は PowerPoint の [ホーム] → [選択] → [オブジェクトの選択と表示]
     で変更できる。

python-pptx では 1 つの段落が複数の run に分割されていることが多く、
``{{price}}`` が ``{{pri`` と ``ce}}`` に割れて素朴な置換では当たらない。
本スクリプトは段落単位でテキストを結合してから置換するため、この問題を踏まない。

依存: python-pptx, Pillow
"""

from __future__ import annotations

import argparse
import io
import json
import os
import re
import sys
from pathlib import Path
from typing import Any, Iterable

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

PLACEHOLDER_RE = re.compile(r"\{\{\s*([A-Za-z0-9_.]+)\s*\}\}")
IMAGE_FRAME_RE = re.compile(r"^IMG:([A-Za-z0-9_]+)$")

NULL_PLACEHOLDER = os.environ.get("PPTX_NULL_PLACEHOLDER", "―")
IMAGE_MAX_PX = int(os.environ.get("PPTX_IMAGE_MAX_PX", "1600"))
JPEG_QUALITY = int(os.environ.get("PPTX_IMAGE_JPEG_QUALITY", "85"))

# Windows の既定コンソール（cp932）で日本語や記号が化けないようにする。
for _stream in (sys.stdout, sys.stderr):
    if hasattr(_stream, "reconfigure"):
        _stream.reconfigure(encoding="utf-8", errors="replace")


# --------------------------------------------------------------------------
# 値の整形
# --------------------------------------------------------------------------

def format_value(value: Any, spec: dict) -> str:
    """DB の生の値を資料に載せる表記へ。単位はテンプレート側に持たせず
    ここで付ける。テンプレートを差し替えても表記が揺れないようにするため。"""
    if value is None or value == "":
        return NULL_PLACEHOLDER

    unit = spec.get("unit")
    if spec.get("type") == "integer" and unit == "円":
        return f"{value:,}円" if value < 10_000_000 else f"{value / 10_000:,.0f}万円"
    if spec.get("type") == "integer" and unit and unit.startswith("円"):
        return f"{value:,}円{unit[1:]}"
    if spec.get("type") == "number" and unit == "㎡":
        return f"{value:,.2f}㎡（{value / 3.305785:,.2f}坪）"
    if spec.get("type") == "number" and unit == "%":
        return f"{value:.2f}%"
    if spec.get("type") == "date_ym" and isinstance(value, str) and len(value) == 7:
        return f"{value[:4]}年{int(value[5:]):d}月"
    if isinstance(value, int):
        return f"{value:,}"
    return str(value)


def format_station(station: dict) -> str:
    parts = [station.get("line") or "", f"{station.get('station') or ''}駅"]
    if station.get("bus_minutes"):
        parts.append(f"バス{station['bus_minutes']}分")
    if station.get("walk_minutes"):
        parts.append(f"徒歩{station['walk_minutes']}分")
    text = " ".join(p for p in parts if p.strip() and p != "駅")
    return text or NULL_PLACEHOLDER


def build_context(data: dict, defs: dict, mark_review: bool) -> dict[str, str]:
    """{{key}} → 差し込む文字列 の対応表を作る。"""
    specs = {f["key"]: f for f in defs["fields"]}
    fields = data.get("fields", {})
    context: dict[str, str] = {}

    for key, spec in specs.items():
        name = spec.get("pptx")
        if not name:
            continue
        env = fields.get(key) or {}
        text = format_value(env.get("value"), spec)
        # 要確認の値をそのまま資料に出すと、確認されないまま客先へ出てしまう。
        # 印刷しても目に入るよう本文側に印を残す。
        if mark_review and env.get("needs_review") and env.get("value") is not None:
            text = f"{text} ※要確認"
        context[name] = text

    stations = data.get("stations", [])
    for index, name in enumerate(defs["repeated_fields"]["stations"]["pptx"]):
        context[name] = format_station(stations[index]) if index < len(stations) else ""

    for meta_key, meta_value in (data.get("meta") or {}).items():
        context[f"meta.{meta_key}"] = "" if meta_value is None else str(meta_value)

    return context


# --------------------------------------------------------------------------
# テキスト差し込み
# --------------------------------------------------------------------------

def _replace_in_text_frame(text_frame, context: dict[str, str], missing: set[str]) -> None:
    """段落単位で置換する。

    run をまたいだプレースホルダに対応するため、段落のテキストを結合して
    置換し、結果を先頭 run に書き戻して残りの run を空にする。書式は
    先頭 run のものに揃う。プレースホルダの途中で書式を変えているテンプレート
    はそもそも意図が読めないので、これで実害はない。
    """
    for paragraph in text_frame.paragraphs:
        runs = paragraph.runs
        if not runs:
            continue
        original = "".join(run.text for run in runs)
        if "{{" not in original:
            continue

        def substitute(match: re.Match) -> str:
            key = match.group(1)
            if key in context:
                return context[key]
            missing.add(key)
            return match.group(0)

        replaced = PLACEHOLDER_RE.sub(substitute, original)
        if replaced == original:
            continue
        runs[0].text = replaced
        for run in runs[1:]:
            run.text = ""


def _walk_shapes(shapes) -> Iterable:
    """グループ図形の中身まで再帰的に辿る。"""
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes)


def fill_text(prs, context: dict[str, str]) -> set[str]:
    missing: set[str] = set()
    for slide in prs.slides:
        for shape in _walk_shapes(slide.shapes):
            if shape.has_text_frame:
                _replace_in_text_frame(shape.text_frame, context, missing)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        _replace_in_text_frame(cell.text_frame, context, missing)
        if slide.has_notes_slide:
            _replace_in_text_frame(slide.notes_slide.notes_text_frame, context, missing)
    return missing


# --------------------------------------------------------------------------
# 画像差し込み
# --------------------------------------------------------------------------

def prepare_image(path: Path) -> tuple[io.BytesIO, int, int]:
    """EXIF の向きを反映し、長辺を上限まで縮めたバイト列を返す。

    スマホ撮影の写真は EXIF Orientation を持っており、そのまま貼ると
    横倒しになる。また元画像は数MBあることが多く、8枚貼ると pptx が
    メールに添付できない大きさになる。
    """
    with Image.open(path) as image:
        image = ImageOps.exif_transpose(image)
        if max(image.size) > IMAGE_MAX_PX:
            image.thumbnail((IMAGE_MAX_PX, IMAGE_MAX_PX), Image.LANCZOS)

        buffer = io.BytesIO()
        if image.mode in ("RGBA", "LA", "P"):
            image.convert("RGBA").save(buffer, format="PNG", optimize=True)
        else:
            image.convert("RGB").save(buffer, format="JPEG",
                                      quality=JPEG_QUALITY, optimize=True)
        width, height = image.size
    buffer.seek(0)
    return buffer, width, height


def fit_box(frame: tuple[int, int, int, int], px_w: int, px_h: int) -> tuple[int, int, int, int]:
    """枠内にアスペクト比を保って収め、中央寄せした位置とサイズ（EMU）。

    間取り図を引き伸ばすと寸法が狂って見えるので、必ず contain で入れる。
    """
    left, top, box_w, box_h = frame
    scale = min(box_w / px_w, box_h / px_h)
    width, height = int(px_w * scale), int(px_h * scale)
    return (left + (box_w - width) // 2, top + (box_h - height) // 2, width, height)


def fill_images(prs, images: list[dict], images_dir: Path,
                drop_empty: bool = False) -> tuple[list[str], list[str]]:
    """名前が IMG:<role> の図形を、対応する画像に置き換える。

    対応する画像が無い枠は既定では残す。テンプレート側で「写真準備中」の
    体裁を作り込んでいることが多く、消すとレイアウトが崩れるため。
    素の四角形をそのまま出したくない場合は drop_empty=True で削除する。
    """
    by_role: dict[str, list[dict]] = {}
    for item in images:
        by_role.setdefault(item.get("role", "other"), []).append(item)

    placed: list[str] = []
    empty_frames: list[str] = []

    for slide in prs.slides:
        for shape in list(_walk_shapes(slide.shapes)):
            match = IMAGE_FRAME_RE.match(shape.name or "")
            if not match:
                continue
            role = match.group(1)
            queue = by_role.get(role, [])
            if not queue:
                empty_frames.append(shape.name)
                if drop_empty:
                    shape._element.getparent().remove(shape._element)
                continue

            item = queue.pop(0)
            source = images_dir / item["file"]
            if not source.exists():
                empty_frames.append(f"{shape.name}（{item['file']} が見つからない）")
                if drop_empty:
                    shape._element.getparent().remove(shape._element)
                continue

            frame = (shape.left, shape.top, shape.width, shape.height)
            buffer, px_w, px_h = prepare_image(source)
            left, top, width, height = fit_box(frame, px_w, px_h)
            slide.shapes.add_picture(buffer, Emu(left), Emu(top), Emu(width), Emu(height))
            shape._element.getparent().remove(shape._element)
            placed.append(f"{shape.name} ← {item['file']}")

    return placed, empty_frames


# --------------------------------------------------------------------------

def main() -> int:
    parser = argparse.ArgumentParser(description="PowerPoint テンプレートに物件情報を差し込む")
    parser.add_argument("--template", type=Path, required=True)
    parser.add_argument("--data", type=Path, required=True,
                        help="extract_property.py / nearest_station.py の出力 JSON")
    parser.add_argument("--fields", type=Path,
                        default=Path(__file__).resolve().parent.parent / "assets" / "property_fields.json")
    parser.add_argument("--images-dir", type=Path, default=Path("."))
    parser.add_argument("--out", type=Path, required=True)
    parser.add_argument("--no-review-marker", action="store_true",
                        help="要確認項目に ※要確認 を付けない")
    parser.add_argument("--drop-empty-frames", action="store_true",
                        help="対応する画像が無い IMG: 枠を削除する（既定は残す）")
    args = parser.parse_args()

    defs = json.loads(args.fields.read_text(encoding="utf-8"))
    data = json.loads(args.data.read_text(encoding="utf-8"))

    prs = Presentation(str(args.template))
    context = build_context(data, defs, mark_review=not args.no_review_marker)

    missing = fill_text(prs, context)
    placed, empty_frames = fill_images(prs, data.get("images", []), args.images_dir,
                                       drop_empty=args.drop_empty_frames)

    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.out))

    print(f"生成: {args.out}")
    for line in placed:
        print(f"  画像配置: {line}")
    if empty_frames:
        print(f"  画像なしの枠: {', '.join(empty_frames)}")
    if missing:
        # テンプレートの誤字や項目定義の追加漏れ。放置すると資料に
        # {{...}} がそのまま印刷されるので必ず気付けるようにする。
        print(f"  [警告] 未定義のプレースホルダ: {', '.join(sorted(missing))}")
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
