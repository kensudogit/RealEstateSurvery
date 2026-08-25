#!/usr/bin/env python3
"""物件資料テンプレートの雛形を生成する。

    python tools/make_starter_templates.py --out templates/

営業が普段使っている pptx をそのまま流用するのが本筋だが、それが揃うまでの
間、パイプラインを端まで通して動作確認するための土台として使う。
生成されるのは 2 種類。

  mysoku_a4.pptx   1 枚もののマイソク（販売図面）
  summary_a4.pptx  2 枚の物件概要書

差し込みの約束ごとは 2 つだけ。

  1. テキストは {{property_name}} のように二重波かっこ。キーは
     config/property_fields.json の pptx の値。
  2. 画像枠は図形の**名前**を IMG:exterior のようにする。図形の位置と
     サイズが画像の枠になる。名前は PowerPoint の
     [ホーム] → [選択] → [オブジェクトの選択と表示] で変更できる。

依存: python-pptx
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

for _stream in (sys.stdout, sys.stderr):
    if hasattr(_stream, "reconfigure"):
        _stream.reconfigure(encoding="utf-8", errors="replace")

A4_WIDTH = Inches(11.69)   # A4 横
A4_HEIGHT = Inches(8.27)

INK = RGBColor(0x1C, 0x20, 0x24)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LINE = RGBColor(0xD5, 0xD9, 0xDF)
FRAME_FILL = RGBColor(0xF0, 0xF2, 0xF5)

FONT = "Yu Gothic"


def _textbox(slide, left, top, width, height, text, *, size=11,
             bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    return box


def _image_frame(slide, name, left, top, width, height, caption):
    """画像の差し込み枠。図形名が差し込み先の識別子になる。

    画像が無い場合この図形は残る（テンプレート側の「写真準備中」の体裁を
    壊さないため）ので、雛形では薄いグレーで塗ってキャプションを入れておく。
    """
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = FRAME_FILL
    shape.line.color.rgb = LINE
    frame = shape.text_frame
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = caption
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED
    run.font.name = FONT
    return shape


def _spec_table(slide, left, top, width, rows: list[tuple[str, str]],
                columns: int = 2) -> None:
    """物件概要の表。実テンプレートも表で組まれていることがほとんど。"""
    per_column = -(-len(rows) // columns)
    column_width = width // columns

    for column_index in range(columns):
        chunk = rows[column_index * per_column:(column_index + 1) * per_column]
        if not chunk:
            continue
        table_shape = slide.shapes.add_table(
            len(chunk), 2,
            left + column_width * column_index, top,
            column_width - Inches(0.15), Inches(0.3) * len(chunk),
        )
        table = table_shape.table
        table.columns[0].width = Emu(int((column_width - Inches(0.15)) * 0.38))
        table.columns[1].width = Emu(int((column_width - Inches(0.15)) * 0.62))

        for row_index, (label, placeholder) in enumerate(chunk):
            for cell_index, text in enumerate((label, placeholder)):
                cell = table.cell(row_index, cell_index)
                cell.text = ""
                paragraph = cell.text_frame.paragraphs[0]
                run = paragraph.add_run()
                run.text = text
                run.font.size = Pt(10)
                run.font.bold = cell_index == 0
                run.font.color.rgb = MUTED if cell_index == 0 else INK
                run.font.name = FONT


def build_mysoku(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = A4_WIDTH
    presentation.slide_height = A4_HEIGHT
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    _textbox(slide, Inches(0.5), Inches(0.35), Inches(7.5), Inches(0.6),
             "{{property_name}}", size=26, bold=True)
    _textbox(slide, Inches(0.5), Inches(1.0), Inches(7.5), Inches(0.3),
             "{{property_type}} / {{deal_type}}", size=11, color=MUTED)

    _textbox(slide, Inches(0.5), Inches(1.45), Inches(4.6), Inches(0.5),
             "{{price}}", size=22, bold=True)
    _textbox(slide, Inches(0.5), Inches(2.0), Inches(4.6), Inches(0.3),
             "表面利回り {{gross_yield}}", size=11, color=MUTED)

    _textbox(slide, Inches(0.5), Inches(2.5), Inches(6.6), Inches(0.3),
             "所在地  {{address_display}}", size=11)
    _textbox(slide, Inches(0.5), Inches(2.85), Inches(6.6), Inches(0.3),
             "交通  {{access_1}}", size=11)
    _textbox(slide, Inches(0.5), Inches(3.15), Inches(6.6), Inches(0.3),
             "　　　{{access_2}}", size=11)

    _spec_table(slide, Inches(0.5), Inches(3.7), Inches(6.9), [
        ("専有面積", "{{exclusive_area_sqm}}"),
        ("間取り", "{{floor_plan}}"),
        ("築年月", "{{built_year_month}}"),
        ("構造", "{{structure}}"),
        ("所在階", "{{floor_of_unit}} / {{floors_total}}"),
        ("管理費", "{{management_fee}}"),
        ("修繕積立金", "{{repair_reserve}}"),
        ("現況", "{{current_status}}"),
        ("引渡時期", "{{delivery_time}}"),
        ("取引態様", "{{transaction_type}}"),
    ])

    _image_frame(slide, "IMG:exterior", Inches(7.7), Inches(0.35),
                 Inches(3.5), Inches(2.4), "外観写真")
    _image_frame(slide, "IMG:floor_plan", Inches(7.7), Inches(2.95),
                 Inches(3.5), Inches(2.6), "間取り図")
    _image_frame(slide, "IMG:map", Inches(7.7), Inches(5.75),
                 Inches(3.5), Inches(1.8), "地図")

    _textbox(slide, Inches(0.5), Inches(6.3), Inches(6.9), Inches(1.0),
             "備考  {{remarks}}", size=9, color=MUTED)
    _textbox(slide, Inches(0.5), Inches(7.6), Inches(10.7), Inches(0.3),
             "{{source_company}}　{{meta.review_status}}", size=8, color=MUTED)

    presentation.save(str(path))


def build_summary(path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = A4_WIDTH
    presentation.slide_height = A4_HEIGHT

    # 1 枚目: 概要
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _textbox(slide, Inches(0.6), Inches(0.5), Inches(10.5), Inches(0.6),
             "物件概要書", size=13, color=MUTED)
    _textbox(slide, Inches(0.6), Inches(0.95), Inches(10.5), Inches(0.7),
             "{{property_name}}", size=28, bold=True)
    _textbox(slide, Inches(0.6), Inches(1.8), Inches(10.5), Inches(0.3),
             "{{address_display}}　|　{{access_1}}", size=12)

    _spec_table(slide, Inches(0.6), Inches(2.5), Inches(10.5), [
        ("物件種別", "{{property_type}}"),
        ("取引種別", "{{deal_type}}"),
        ("価格", "{{price}}"),
        ("表面利回り", "{{gross_yield}}"),
        ("土地面積", "{{land_area_sqm}}"),
        ("建物面積", "{{building_area_sqm}}"),
        ("専有面積", "{{exclusive_area_sqm}}"),
        ("間取り", "{{floor_plan}}"),
        ("築年月", "{{built_year_month}}"),
        ("構造", "{{structure}}"),
        ("総戸数", "{{units_total}}"),
        ("所在階 / 総階数", "{{floor_of_unit}} / {{floors_total}}"),
        ("現況", "{{current_status}}"),
        ("引渡時期", "{{delivery_time}}"),
        ("取引態様", "{{transaction_type}}"),
        ("元付会社", "{{source_company}}"),
    ], columns=2)

    _textbox(slide, Inches(0.6), Inches(7.6), Inches(10.5), Inches(0.3),
             "{{meta.review_status}}　要確認項目: {{meta.review_fields}}",
             size=8, color=MUTED)

    # 2 枚目: 写真・間取り
    photos = presentation.slides.add_slide(presentation.slide_layouts[6])
    _textbox(photos, Inches(0.6), Inches(0.5), Inches(10.5), Inches(0.5),
             "{{property_name}}　写真・間取り", size=16, bold=True)
    _image_frame(photos, "IMG:exterior", Inches(0.6), Inches(1.2),
                 Inches(5.1), Inches(3.2), "外観写真")
    _image_frame(photos, "IMG:interior", Inches(6.0), Inches(1.2),
                 Inches(5.1), Inches(3.2), "室内写真")
    _image_frame(photos, "IMG:floor_plan", Inches(0.6), Inches(4.6),
                 Inches(5.1), Inches(3.2), "間取り図")
    _image_frame(photos, "IMG:map", Inches(6.0), Inches(4.6),
                 Inches(5.1), Inches(3.2), "地図")

    presentation.save(str(path))


def main() -> int:
    parser = argparse.ArgumentParser(description="物件資料テンプレートの雛形を作る")
    parser.add_argument("--out", type=Path, default=Path("templates"))
    parser.add_argument("--fields", type=Path, default=Path("config/property_fields.json"),
                        help="使用したプレースホルダの検証に使う")
    args = parser.parse_args()

    args.out.mkdir(parents=True, exist_ok=True)
    build_mysoku(args.out / "mysoku_a4.pptx")
    build_summary(args.out / "summary_a4.pptx")

    print(f"生成しました: {args.out / 'mysoku_a4.pptx'}")
    print(f"生成しました: {args.out / 'summary_a4.pptx'}")

    if args.fields.exists():
        definitions = json.loads(args.fields.read_text(encoding="utf-8"))
        available = {f["pptx"] for f in definitions["fields"] if f.get("pptx")}
        available |= set(definitions["repeated_fields"]["stations"]["pptx"])
        print(f"利用可能なプレースホルダ: {len(available)} 個")
        print("  " + "、".join(sorted(available)))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
