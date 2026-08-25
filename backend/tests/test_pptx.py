"""資料生成のテスト。

run 分割されたプレースホルダと画像の contain フィットが実際に効くかを、
テンプレートを組み立てて確認する。ここが壊れると {{price}} がそのまま
印刷された資料が客先に出る。
"""

from __future__ import annotations

import pytest

pptx = pytest.importorskip("pptx")
PIL = pytest.importorskip("PIL")

from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from app.services.pptx import (  # noqa: E402
    build_context, fill_images, fill_text, fit_box, format_station, format_value,
)


@pytest.fixture
def sample_data():
    def envelope(value, review=False, confidence=0.95):
        return {"value": value, "confidence": confidence, "evidence": "原文",
                "needs_review": review, "review_reasons": []}

    return {
        "fields": {
            "property_name": envelope("渋谷ハイツ"),
            "price": envelope(48_000_000, review=True),
            "exclusive_area_sqm": envelope(62.5),
            "built_year_month": envelope("1993-03"),
            "floor_plan": envelope(None),
        },
        "stations": [
            {"line": "JR山手線", "station": "渋谷", "walk_minutes": 7,
             "distance_m": 520, "bus_minutes": None, "source": "extracted"},
        ],
        "images": [],
        "meta": {"review_status": "要確認"},
    }


@pytest.fixture
def template(tmp_path):
    """run が分割されたプレースホルダを含むテンプレートを作る。

    PowerPoint 上で一度でも編集すると {{price}} が {{pri と ce}} に割れる。
    実テンプレートで最も多いつまずき。
    """
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(6), Inches(0.6))
    paragraph = box.text_frame.paragraphs[0]
    for chunk in ["{{prop", "erty_", "name}}", " / ", "{{access_1}}"]:
        run = paragraph.add_run()
        run.text = chunk
        run.font.size = Pt(20)

    table = slide.shapes.add_table(3, 2, Inches(0.5), Inches(1.2), Inches(4), Inches(1.5)).table
    for index, (label, placeholder) in enumerate([
        ("価格", "{{price}}"),
        ("専有面積", "{{exclusive_area_sqm}}"),
        ("間取り", "{{floor_plan}}"),
    ]):
        table.cell(index, 0).text = label
        table.cell(index, 1).text = placeholder

    typo = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(4), Inches(0.4))
    typo.text_frame.paragraphs[0].add_run().text = "{{proprety_name}}"

    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(1.2),
                                   Inches(3.0), Inches(2.5))
    frame.name = "IMG:exterior"

    path = tmp_path / "template.pptx"
    presentation.save(str(path))
    return path


class TestFormatValue:
    @pytest.mark.parametrize(
        "value,spec,expected",
        [
            (48_000_000, {"type": "integer", "unit": "円"}, "4,800万円"),
            (100_000_000, {"type": "integer", "unit": "円"}, "1億円"),
            (128_000_000, {"type": "integer", "unit": "円"}, "1億2,800万円"),
            (8_500_000, {"type": "integer", "unit": "円"}, "8,500,000円"),
            # 万で割り切れない端数は丸めずに円で出す（競売・持分など）
            (48_123_456, {"type": "integer", "unit": "円"}, "48,123,456円"),
            (12_000, {"type": "integer", "unit": "円/月"}, "12,000円/月"),
            (7.2, {"type": "number", "unit": "%"}, "7.20%"),
            ("1993-03", {"type": "date_ym"}, "1993年3月"),
            (None, {"type": "string"}, "―"),
        ],
    )
    def test_format(self, value, spec, expected):
        assert format_value(value, spec) == expected

    def test_sqm_includes_tsubo(self):
        assert format_value(62.5, {"type": "number", "unit": "㎡"}) == "62.50㎡（18.91坪）"

    def test_null_is_not_blank(self):
        """空文字にすると「未取得」と「もともと空欄」が読み手に区別できない。"""
        assert format_value(None, {"type": "string"}) == "―"


def test_format_station():
    assert format_station(
        {"line": "JR山手線", "station": "渋谷", "walk_minutes": 7, "bus_minutes": None}
    ) == "JR山手線 渋谷駅 徒歩7分"


class TestContext:
    def test_review_marker_is_added(self, sample_data):
        context = build_context(sample_data, mark_review=True)
        assert context["price"].endswith("※要確認")

    def test_review_marker_can_be_suppressed(self, sample_data):
        context = build_context(sample_data, mark_review=False)
        assert "要確認" not in context["price"]

    def test_missing_station_slot_is_empty(self, sample_data):
        context = build_context(sample_data)
        assert context["access_1"] == "JR山手線 渋谷駅 徒歩7分"
        assert context["access_2"] == ""


class TestFillText:
    def test_split_runs_are_replaced(self, template, sample_data):
        presentation = Presentation(str(template))
        missing = fill_text(presentation, build_context(sample_data))

        texts = [
            shape.text_frame.text
            for shape in presentation.slides[0].shapes
            if shape.has_text_frame
        ]
        assert any("渋谷ハイツ / JR山手線 渋谷駅 徒歩7分" == text for text in texts)
        # テンプレートの誤字は握りつぶさず報告する
        assert missing == {"proprety_name"}

    def test_table_cells_are_replaced(self, template, sample_data):
        presentation = Presentation(str(template))
        fill_text(presentation, build_context(sample_data))

        table = next(s.table for s in presentation.slides[0].shapes if s.has_table)
        rows = {table.cell(i, 0).text: table.cell(i, 1).text for i in range(3)}
        assert rows["価格"] == "4,800万円 ※要確認"
        assert rows["専有面積"] == "62.50㎡（18.91坪）"
        assert rows["間取り"] == "―"


class TestFitBox:
    def test_contain_and_center(self):
        # 3:1 の横長画像を 3:2.5 の枠に入れる
        left, top, width, height = fit_box((100, 200, 3000, 2500), 1200, 400)
        assert width == 3000
        assert height == 1000
        assert left == 100
        assert top == 200 + (2500 - 1000) // 2

    def test_aspect_ratio_is_preserved(self):
        _, _, width, height = fit_box((0, 0, 1000, 1000), 400, 200)
        assert width / height == pytest.approx(2.0)


class TestFillImages:
    def test_image_replaces_named_frame(self, template, tmp_path, sample_data):
        photo = tmp_path / "gaikan.jpg"
        Image.new("RGB", (1200, 400), (200, 120, 90)).save(photo)

        presentation = Presentation(str(template))
        placed, empty = fill_images(
            presentation,
            [{"file": "gaikan.jpg", "role": "exterior", "storage_path": str(photo)}],
        )

        assert len(placed) == 1
        assert empty == []
        remaining = [s.name for s in presentation.slides[0].shapes if s.name.startswith("IMG:")]
        assert remaining == []

    def test_empty_frame_is_kept_by_default(self, template):
        """テンプレート側で「写真準備中」の体裁を作り込んでいることが多く、
        消すとレイアウトが崩れる。"""
        presentation = Presentation(str(template))
        placed, empty = fill_images(presentation, [])

        assert placed == []
        assert empty == ["IMG:exterior"]
        remaining = [s.name for s in presentation.slides[0].shapes if s.name.startswith("IMG:")]
        assert remaining == ["IMG:exterior"]

    def test_empty_frame_can_be_dropped(self, template):
        presentation = Presentation(str(template))
        fill_images(presentation, [], drop_empty=True)
        remaining = [s.name for s in presentation.slides[0].shapes if s.name.startswith("IMG:")]
        assert remaining == []
